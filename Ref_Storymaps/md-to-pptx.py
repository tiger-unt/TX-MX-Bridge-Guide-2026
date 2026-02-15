"""
Generate PowerPoint presentation from StoryMap Comparison Analysis

Usage:
    python create_comparison_pptx.py

Requirements:
    pip install python-pptx

Output:
    StoryMap_Comparison_Analysis.pptx (in same directory)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# Color scheme - TxDOT-inspired
DARK_BLUE = RGBColor(0, 51, 102)
LIGHT_BLUE = RGBColor(0, 102, 153)
ACCENT_ORANGE = RGBColor(204, 102, 0)
DARK_GRAY = RGBColor(64, 64, 64)
LIGHT_GRAY = RGBColor(240, 240, 240)
WHITE = RGBColor(255, 255, 255)
LIGHT_TEXT = RGBColor(200, 200, 200)

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = LIGHT_TEXT
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, title):
    """Add a section divider slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Accent bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3), prs.slide_width, Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, two_column=False, col1_title="", col2_title="", col1_bullets=None, col2_bullets=None):
    """Add a content slide with bullets"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    if two_column and col1_bullets and col2_bullets:
        # Two column layout
        # Column 1 header
        col1_header = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(6), Inches(0.5))
        tf = col1_header.text_frame
        p = tf.paragraphs[0]
        p.text = col1_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = LIGHT_BLUE
        
        # Column 1 content
        col1_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(6), Inches(5))
        tf = col1_box.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(col1_bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + bullet
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(6)
        
        # Column 2 header
        col2_header = slide.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(6), Inches(0.5))
        tf = col2_header.text_frame
        p = tf.paragraphs[0]
        p.text = col2_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = ACCENT_ORANGE
        
        # Column 2 content
        col2_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.9), Inches(6), Inches(5))
        tf = col2_box.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(col2_bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + bullet
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(6)
    else:
        # Single column
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + bullet
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(8)
    
    return slide

def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Table
    num_cols = len(headers)
    num_rows = len(rows) + 1
    table_width = Inches(12.333)
    table_height = Inches(5.5)
    
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.4), table_width, table_height).table
    
    # Set column widths
    col_width = table_width / num_cols
    for i in range(num_cols):
        table.columns[i].width = int(col_width)
    
    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GRAY
            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
    
    return slide


def create_presentation():
    """Build the complete presentation"""
    # Create presentation with widescreen dimensions
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ============================================
    # BUILD THE PRESENTATION
    # ============================================

    # Slide 1: Title
    add_title_slide(prs, 
        "Texas-Mexico Border Story Maps",
        "Comparison Analysis")

    # Slide 2: Documents Analyzed
    add_content_slide(prs, "Story Maps Analyzed", [
        "Connectivity Plan Story Map: Texas-Mexico Border Region Connectivity Plan",
        "Border Crossings Guide Story Map: Texas-Mexico International Border Crossings",
        "",
        "Both published by TxDOT Transportation Planning and Programming Division",
        "Both use Bureau of Transportation Statistics (BTS) Transborder Data"
    ])

    # Slide 3: Key Findings
    add_content_slide(prs, "Key Findings", [
        "Primary Difference: Regional vs. Individual Crossing Focus",
        "",
        "Connectivity Plan: Defines binational economic regions (West, Central, South) spanning 100 miles on both sides of the border",
        "",
        "Border Crossings Guide: Summary and detailed port of entry (trade) and border crossing/bridge information by TxDOT district: El Paso, Laredo, Rio Grande Valley"
    ])

    # Slide 4: Data Consistency
    add_content_slide(prs, "Data Consistency Between Documents", [
        "Both StoryMaps use Bureau of Transportation Statistics (BTS) and are consistent in identifying:",
        "",
        "Top commodity types: machinery/electrical equipment, transportation equipment, metals, plastics/rubbers",
        "",
        "Truck dominance: ~80% of trade value crosses by truck",
        "",
        "Similar pandemic recovery patterns (2020-2021 disruptions, subsequent recovery)"
    ])

    # Slide 5: Regional Definition Comparison Table
    add_table_slide(prs, "Regional Definition Implications", 
        ["Aspect", "Connectivity Plan", "Border Crossings Guide"],
        [
            ["Geographic Scope", "Binational (100-mile buffer both sides)", "TX-MX bridge crossings by TxDOT District"],
            ["Population", "13.6 million (both countries)", "Not included"],
            ["Boundary Type", "Economic/functional regions", "TxDOT border district"],
            ["Mexican Info", "Socio-economic data for Mexican states", "Bridge info (owner, operator, tolls)"],
            ["Planning Horizon", "Long-range (2050 projections)", "Focus on 2024 data"],
            ["Key Metrics", "Employment, GDP, freight tonnage", "Trade value/volumes by POE/bridge"]
        ])

    # Slide 6: Content Scope - Trade & Commerce
    add_table_slide(prs, "Content Scope: Trade & Commerce",
        ["Content Category", "Connectivity Plan", "Border Crossings Guide"],
        [
            ["Trade Value Data", "✓ (2023, by region)", "✓ (2024, by Port of Entry)"],
            ["Crossing Volume Data", "Limited", "✓ Comprehensive (all modes)"],
            ["Economic Projections", "✓ (2050)", "✗"],
            ["Truck Freight Tonnage", "✓", "✗"],
            ["Commodity Types", "✓", "✓"],
            ["Seasonal Patterns", "✓ (monthly variations)", "✗"],
            ["Historical Trends", "✓ (2019-2023, 5 years)", "✓ (2014-2024, 10 years)"]
        ])

    # Slide 7: Content Scope - Infrastructure
    add_table_slide(prs, "Content Scope: Infrastructure",
        ["Content Category", "Connectivity Plan", "Border Crossings Guide"],
        [
            ["Bridge/Crossing Names", "Limited", "✓ Comprehensive (34 bridges)"],
            ["Individual Crossing Info", "✗", "✓"],
            ["Highway Corridors", "✓ (detailed analysis)", "Reference bridge access roads"],
            ["Rail Infrastructure", "✓ (4 crossings)", "✓ (detailed, individual)"],
            ["Airports", "✓ (8 commercial)", "Lists airports in POEs"],
            ["Maritime Ports", "✓ (2 deep, 3 shallow)", "Mentioned Port of Brownsville"]
        ])

    # Slide 8: Content Scope - Operations & Planning
    add_table_slide(prs, "Content Scope: Operations & Planning",
        ["Content Category", "Connectivity Plan", "Border Crossings Guide"],
        [
            ["Pedestrian Volumes", "✗", "✓ (detailed by crossing)"],
            ["Passenger Vehicle Volumes", "✗", "✓ (detailed by crossing)"],
            ["Commercial Truck Volumes", "✓ (regional)", "✓ (by crossing)"],
            ["Processing Programs (SENTRI, FAST)", "✗", "✓"],
            ["Employment Data/Projections", "✓ (2022, 2050)", "✗"],
            ["GDP Data", "✓", "✗"],
            ["Corridor Needs Assessment", "✓ (scoring methodology)", "✗"]
        ])

    # Slide 9: Trade Statistics Reconciliation
    add_table_slide(prs, "Trade Statistics Reconciliation",
        ["Metric", "Connectivity Plan", "Border Crossings Guide", "Notes"],
        [
            ["Data Year", "2023", "2024", "Different reporting periods"],
            ["Total Trade", "$468B", "$552.9B", "Year difference + methodology"],
            ["Truck Trade", "~80% of total", "$454.5B (82%)", "Consistent proportion"],
            ["Rail Trade", "~16% of total", "$88.8B (16%)", "Exact match"],
            ["Other Modes", "~4%", "$9.6B (2%)", "Pipeline, air, vessel, FTZ"]
        ])

    # Slide 10: Infrastructure Inventory
    add_table_slide(prs, "Infrastructure Inventory Reconciliation",
        ["Category", "Connectivity Plan", "Border Crossings Guide"],
        [
            ["Border Crossings", "28 total", "34 operational"],
            ["Truck-Capable", "18", "14 (4+4+6)"],
            ["Rail Bridges", "4 operational", "5 operational"],
            ["Ports of Entry", "Not specified", "12"],
            ["Airports", "8 commercial", "Lists in POEs"],
            ["Deep Draft Ports", "2", "1 (Brownsville)"]
        ])

    # Slide 11: Section - Unique Value
    add_section_slide(prs, "Unique Value Propositions")

    # Slide 12: Connectivity Plan Strengths
    add_content_slide(prs, "Connectivity Plan: Strategic Planning Document",
        two_column=True,
        col1_title="Core Strengths",
        col1_bullets=[
            "Binational Economic Perspective - integrates both sides of border",
            "Long-Range Planning Framework - 2050 projections",
            "Multimodal Connectivity Analysis - First/Last Mile",
            "Granular Freight Flow Data - county-to-county",
            "Infrastructure Needs Assessment - 661 projects",
            "Climate Resilience Focus - extreme weather events",
            "Policy Implementation Tool - 22 policies, 153 programs"
        ],
        col2_title="Best Used For",
        col2_bullets=[
            "Long-range transportation planning",
            "Capital investment prioritization",
            "Economic development strategy",
            "Binational coordination",
            "Infrastructure needs assessment",
            "Grant applications requiring economic impact",
            "Environmental and climate resilience planning"
        ],
        bullets=[])

    # Slide 13: Border Crossings Guide Strengths
    add_content_slide(prs, "Border Crossings Guide: Operational Inventory",
        two_column=True,
        col1_title="Core Strengths",
        col1_bullets=[
            "Comprehensive Border Crossing Inventory - all 34 crossings",
            "Border Crossing Performance Data - 10-year trends",
            "Practical Crossing Information - SENTRI, FAST, Ready Lane",
            "Trade and Volume Data at Two Levels - district & crossing",
            "Future Infrastructure Projects - 10 Presidential Permits",
            "Traveler-Friendly Format - interactive carousels",
            "Regular Update Cycle - annual publication"
        ],
        col2_title="Best Used For",
        col2_bullets=[
            "Comparing districts and individual crossings",
            "Selecting which crossing to use (traveler planning)",
            "Understanding crossing capabilities by mode",
            "District-level and crossing-level reporting",
            "Future crossing projects tracking",
            "Media fact-checking on crossing statistics",
            "Academic research on border operations"
        ],
        bullets=[])

    # Slide 14: Methodology Differences
    add_content_slide(prs, "Methodology Differences Explained", [
        "Regional vs. Crossing-Based Organization:",
        "  - Connectivity Plan: functional economic regions for corridor analysis",
        "  - Border Crossings Guide: TxDOT districts as administrative grouping",
        "",
        "Time Period Differences:",
        "  - Connectivity Plan: 2023 data (published 2025)",
        "  - Border Crossings Guide: 2024 data (annual updates)",
        "",
        "Infrastructure Counting:",
        "  - Connectivity Plan: counts POEs/major facilities (28)",
        "  - Border Crossings Guide: counts individual crossing points (34)"
    ])

    # Slide 15: Recommendations
    add_content_slide(prs, "Recommendations", [
        "Maintain both documents with their distinct focuses",
        "",
        "Do NOT attempt to force geographic alignment between them",
        "",
        "Link from Connectivity Plan regional analysis to relevant Border Crossings Guide POE pages",
        "",
        "Link from Border Crossings Guide crossing data to relevant Connectivity Plan corridor needs",
        "",
        "Create brief guide: 'Use Connectivity Plan for corridor planning; use Border Crossings Guide for crossing selection and performance'"
    ])

    # Slide 16: Conclusion
    add_content_slide(prs, "Conclusion", [
        "The two Story Maps are COMPLEMENTARY resources serving different purposes:",
        "",
        "Connectivity Plan: Strategic, binational, long-range planning analyzing economic regions and transportation corridors",
        "",
        "Border Crossings Guide: Operational data for border crossings with district totals and individual crossing performance",
        "",
        "When used together, they provide complete picture:",
        "  - Regional corridor planning + crossing-level operational data",
        "  - Long-range infrastructure needs (2050) + current performance (2024)",
        "  - Binational economic analysis + U.S.-side crossing operations"
    ])

    return prs


if __name__ == "__main__":
    # Create the presentation
    prs = create_presentation()
    
    # Save in the same directory as this script
    script_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(script_dir, "StoryMap_Comparison_Analysis.pptx")
    
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
